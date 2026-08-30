# =============================================================================
# DBV PDF2Deck — Local OCR · Visual Canvas · PPTX Export
# Copyright (c) 2026 David Bueno Vallejo · https://davidbuenov.com
# Licensed under the MIT License. See LICENSE for details.
# =============================================================================
"""
Motor de re-ensamblaje y exportación hacia ecosistema Office (PPTX) y estándar PDF.
Dependencias principales: python-pptx, reportlab, pypdf
"""
from __future__ import annotations

import base64
import io
import logging
import tempfile
import zipfile
from pathlib import Path
from typing import Any, cast

from pypdf import PdfReader, PdfWriter
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfgen.canvas import Canvas

from .result import Err, Ok
from .settings import EXPORT_JPEG_QUALITY, EXPORT_PNG_MAX_BYTES, resolve_export_dpi

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Pt
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    Pt = None
    Emu = None
    RGBColor = None
    PP_ALIGN = None
    MSO_ANCHOR = None
    PPTX_AVAILABLE = False


logger = logging.getLogger(__name__)


ALIGN_MAP_PDF = {"left": 0, "center": 1, "right": 2}
ALIGN_MAP_PPTX = {"left": None, "center": None, "right": None}  # se rellena tras import
EMU_PER_PT = 12700


def _decode_image(b64_string: str) -> bytes:
    """Extrae la trama binaria en crudo eliminando cabeceras Web."""
    if "," in b64_string:
        b64_string = b64_string.split(",")[1]
    return base64.b64decode(b64_string)


def _hex_to_rgb(hex_code: str) -> tuple[int, int, int]:
    """Parseador simple Hex a RGB de 3 Tuplas enteras para PPTX/PDF"""
    h = hex_code.lstrip('#')
    if len(h) != 6:
        return (255, 255, 255) # Fallback Blanco puro
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _to_pdf_font(font_name: str | None, is_bold: bool = False, is_italic: bool = False) -> str:
    """
    Mapea nombres de fuentes variadas hacia Base-14 de PDF con soporte para Bold e Italic.
    Retorna nombres estándar de PDF Base-14.
    """
    if not font_name:
        base_name = "Helvetica"
    else:
        name = font_name.lower()
        if "cour" in name:
            base_name = "Courier"
        elif "times" in name or "georgia" in name or "serif" in name:
            base_name = "Times-Roman"
        else:
            base_name = "Helvetica"
    
    # Construir nombre con sufijos de estilo según estándar PDF Base-14
    if is_bold and is_italic:
        if base_name.startswith("Times"):
            return "Times-BoldItalic"
        elif base_name.startswith("Courier"):
            return "Courier-BoldOblique"
        else:  # Helvetica
            return "Helvetica-BoldOblique"
    elif is_bold:
        if base_name.startswith("Times"):
            return "Times-Bold"
        elif base_name.startswith("Courier"):
            return "Courier-Bold"
        else:  # Helvetica
            return "Helvetica-Bold"
    elif is_italic:
        if base_name.startswith("Times"):
            return "Times-Italic"
        elif base_name.startswith("Courier"):
            return "Courier-Oblique"
        else:  # Helvetica
            return "Helvetica-Oblique"
    else:
        return base_name


def _page_scale_to_pdf_points(page_data: dict) -> tuple[float, float]:
    render_w = float(page_data.get("render_width_px") or 1.0)
    render_h = float(page_data.get("render_height_px") or 1.0)
    page_w = float(page_data.get("page_width_pt") or render_w)
    page_h = float(page_data.get("page_height_pt") or render_h)

    return page_w / render_w, page_h / render_h


def _page_scale_to_ppt_points(page_data: dict, width_px: float, height_px: float) -> tuple[float, float]:
    """Convierte coordenadas del canvas (px) a puntos tipográficos para PPTX según metadatos de página."""
    page_w_pt = float(page_data.get("page_width_pt") or width_px)
    page_h_pt = float(page_data.get("page_height_pt") or height_px)
    safe_w = max(1.0, float(width_px or 1.0))
    safe_h = max(1.0, float(height_px or 1.0))
    return page_w_pt / safe_w, page_h_pt / safe_h


def _safe_line_spacing(raw_value: float | str | None) -> float:
    try:
        parsed = float(raw_value) if raw_value is not None else 1.15
    except (TypeError, ValueError):
        parsed = 1.15
    return max(0.8, min(3.0, parsed))


def _reportlab_text_lines(text: str, font_name: str, font_size: float, max_width: float) -> list[str]:
    lines: list[str] = []
    for raw_line in str(text or "").split("\n"):
        current = ""
        for word in raw_line.split(" "):
            candidate = f"{current} {word}".strip() if current else word
            if current and pdfmetrics.stringWidth(candidate, font_name, font_size) > max_width:
                lines.append(current)
                current = word
            else:
                current = candidate
        lines.append(current)
    return lines or [""]


def _fit_reportlab_font_size(
    text: str,
    font_name: str,
    requested_size: float,
    width: float,
    height: float,
    line_spacing: float,
) -> tuple[float, list[str], list[str]]:
    size = max(6.0, requested_size)
    minimum_size = 6.0
    while size > minimum_size:
        lines = _reportlab_text_lines(text, font_name, size, width)
        if len(lines) * size * line_spacing <= height:
            return size, lines, []
        size = max(minimum_size, size * 0.9)
    lines = _reportlab_text_lines(text, font_name, size, width)
    capacity = max(1, int(height // (size * line_spacing)))
    return size, lines[:capacity], lines[capacity:]


def _draw_reportlab_textbox(
    canvas: Canvas,
    page_height: float,
    rect: tuple[float, float, float, float],
    text: str,
    font_size: float,
    font_name: str,
    color: tuple[float, float, float],
    align: int,
    line_spacing: float,
    underline: bool,
    label: str,
) -> None:
    x0, y0, x1, y1 = rect
    width = max(1.0, x1 - x0)
    height = max(1.0, y1 - y0)
    actual_size, lines, overflow = _fit_reportlab_font_size(
        text, font_name, font_size, width, height, line_spacing
    )
    if overflow:
        logger.warning("Texto truncado en %s: %d líneas sobrantes", label, len(overflow))
    line_height = actual_size * line_spacing
    canvas.setFont(font_name, actual_size)
    canvas.setFillColorRGB(*color)
    canvas.setStrokeColorRGB(*color)
    for index, line in enumerate(lines):
        text_width = pdfmetrics.stringWidth(line, font_name, actual_size)
        if align == 1:
            x = x0 + (width - text_width) / 2.0
        elif align == 2:
            x = x1 - text_width
        else:
            x = x0
        baseline = page_height - y0 - ((index + 1) * line_height)
        canvas.drawString(x, baseline, line)
        if underline and line:
            underline_y = baseline - max(0.6, actual_size * 0.08)
            canvas.setLineWidth(max(0.6, actual_size * 0.05))
            canvas.line(x, underline_y, x + text_width, underline_y)


def _draw_reportlab_blocks(canvas: Canvas, page_data: dict, page_width: float, page_height: float) -> None:
    scale_x, scale_y = _page_scale_to_pdf_points(page_data)
    export_mode = page_data.get("_export_mode", "only_modified")
    for block_index, block in enumerate(page_data.get("blocks", [])):
        if export_mode == "only_modified" and not bool(block.get("is_modified", False)):
            continue
        x0, y0, x1, y1 = (float(value) for value in block.get("bbox", [0, 0, 0, 0]))
        rect = (x0 * scale_x, y0 * scale_y, x1 * scale_x, y1 * scale_y)
        if not block.get("bg_transparent", False):
            red, green, blue = _hex_to_rgb(str(block.get("bg_color", "#ffffff")))
            canvas.setFillColorRGB(red / 255.0, green / 255.0, blue / 255.0)
            canvas.rect(rect[0], page_height - rect[3], rect[2] - rect[0], rect[3] - rect[1], stroke=0, fill=1)
        red, green, blue = _hex_to_rgb(str(block.get("text_color", "#000000")))
        font_name = _to_pdf_font(block.get("font_family"), bool(block.get("is_bold")), bool(block.get("is_italic")))
        _draw_reportlab_textbox(
            canvas,
            page_height,
            rect,
            str(block.get("text", "")),
            max(6.0, float(block.get("font_size") or 16.0) * scale_y),
            font_name,
            (red / 255.0, green / 255.0, blue / 255.0),
            ALIGN_MAP_PDF.get(str(block.get("text_align", "left")), 0),
            _safe_line_spacing(block.get("line_spacing", 1.15)),
            bool(block.get("is_underline", False)),
            f"página {page_data.get('page_num', 0) + 1}, bloque {block_index + 1}",
        )


def _reportlab_signature(writer: PdfWriter) -> None:
    writer.add_metadata({
        "/Creator": "Generador por davidbuenov.com",
        "/Producer": "DBVPDFEditor",
        "/Subject": "Documento generado con DBVPDFEditor",
        "/Keywords": "DBVPDFEditor, davidbuenov.com",
    })


def _reportlab_image_overlay(page_data: dict, page_width: float, page_height: float, include_image: bool) -> bytes:
    buffer = io.BytesIO()
    canvas = Canvas(buffer, pagesize=(page_width, page_height))
    if include_image and page_data.get("image_base64"):
        canvas.drawImage(
            ImageReader(io.BytesIO(_decode_image(str(page_data["image_base64"])))),
            0,
            0,
            width=page_width,
            height=page_height,
            preserveAspectRatio=True,
            anchor="sw",
        )
    _draw_reportlab_blocks(canvas, page_data, page_width, page_height)
    canvas.showPage()
    canvas.save()
    return buffer.getvalue()


def _build_pdf_export_reportlab(payload: dict) -> bytes:
    writer = PdfWriter()
    for page_data in payload.get("pages", []):
        image_bytes = _decode_image(str(page_data.get("image_base64", "")))
        from PIL import Image
        image = Image.open(io.BytesIO(image_bytes))
        page_width = float(page_data.get("page_width_pt") or image.width)
        page_height = float(page_data.get("page_height_pt") or image.height)
        page_data = {**page_data, "_export_mode": payload.get("export_mode", "only_modified")}
        overlay_bytes = _reportlab_image_overlay(page_data, page_width, page_height, True)
        overlay = PdfReader(io.BytesIO(overlay_bytes))
        if overlay.pages:
            writer.add_page(overlay.pages[0])
    _reportlab_signature(writer)
    output = io.BytesIO()
    writer.write(output)
    return output.getvalue()


def _build_pdf_export_from_original_reportlab(payload: dict, source_pdf_path: Path) -> bytes:
    reader = PdfReader(str(source_pdf_path))
    writer = PdfWriter()
    export_mode = payload.get("export_mode", "only_modified")
    pages_by_number = {
        int(page.get("page_num", 0)): {**page, "_export_mode": export_mode}
        for page in payload.get("pages", [])
    }
    for page_number, source_page in enumerate(reader.pages):
        page_data = pages_by_number.get(page_number)
        if page_data is not None and page_data.get("blocks"):
            has_modifications = any(bool(b.get("is_modified")) for b in page_data.get("blocks", []))
            has_cleaned_bg = bool(page_data.get("ai_cleaned_bg"))

            # Si no hay modificaciones ni fondo IA nuevo, preservamos la página intacta
            if export_mode == "only_modified" and not has_modifications and not has_cleaned_bg:
                writer.add_page(source_page)
                continue

            page_width = float(source_page.mediabox.width)
            page_height = float(source_page.mediabox.height)
            overlay_bytes = _reportlab_image_overlay(
                page_data,
                page_width,
                page_height,
                has_cleaned_bg,
            )
            overlay_reader = PdfReader(io.BytesIO(overlay_bytes))
            if overlay_reader.pages:
                source_page.merge_page(overlay_reader.pages[0])
        writer.add_page(source_page)
    _reportlab_signature(writer)
    output = io.BytesIO()
    writer.write(output)
    return output.getvalue()


def build_pdf_export(payload: dict) -> bytes:
    """
    Restaura la estructura original repintando las alteraciones visuales sobre el lienzo virtual
    cerrándolo nativamente en formato .PDF mediante ReportLab.
    """
    return _build_pdf_export_reportlab(payload)


def build_pdf_export_from_original(payload: dict, source_pdf_path: Path) -> bytes:
    """
    Estrategia unificada para exportar PDF con ediciones preservando calidad máxima
    en páginas no tocadas mediante overlay con ReportLab y pypdf.
    """
    return _build_pdf_export_from_original_reportlab(payload, source_pdf_path)



def _encode_background(image: Any) -> bytes:
    """
    Serializa el fondo de un slide priorizando PNG sin pérdida.

    Si el PNG se dispara (páginas fotográficas a 400-600 DPI dan cientos de MB
    por deck) se pasa a JPEG de alta calidad: sigue siendo muchísimo mejor que
    el raster de 100 DPI del lienzo, que es el problema que se quiere resolver.

    Args:
        image (Image.Image): Página ya rasterizada en RGB.

    Returns:
        bytes: Contenido del fichero de imagen listo para incrustar.
    """
    png_buffer = io.BytesIO()
    image.save(png_buffer, format="PNG")
    png_bytes = png_buffer.getvalue()
    if len(png_bytes) <= EXPORT_PNG_MAX_BYTES:
        return png_bytes

    jpeg_buffer = io.BytesIO()
    image.save(jpeg_buffer, format="JPEG", quality=EXPORT_JPEG_QUALITY, optimize=True)
    jpeg_bytes = jpeg_buffer.getvalue()
    logger.info(
        "Fondo de %.1f MB en PNG sustituido por JPEG q%d de %.1f MB.",
        len(png_bytes) / 1e6, EXPORT_JPEG_QUALITY, len(jpeg_bytes) / 1e6
    )
    return jpeg_bytes


def _high_res_background(page_data: dict, source_pdf_path: Path | None, export_dpi: int) -> bytes | None:
    """
    Devuelve el fondo del slide re-rasterizado del PDF original, o `None`.

    Se descarta en dos casos legítimos: cuando no hay PDF de origen (la entrada
    era una imagen, donde el raster del lienzo ya son los píxeles originales y
    no hay nada que recuperar) y cuando la página se limpió con la goma o el
    inpaint, porque esos píxeles corregidos solo existen en el raster del
    lienzo y volver al original los perdería.

    Args:
        page_data (dict): Página tal y como la envía el frontend.
        source_pdf_path (Path | None): PDF original conservado hasta la exportación.
        export_dpi (int): Resolución solicitada.

    Returns:
        bytes | None: Imagen codificada, o `None` si toca usar el lienzo.
    """
    if source_pdf_path is None or not source_pdf_path.exists():
        return None
    if page_data.get("ai_cleaned_bg"):
        return None

    from .pdf_renderer import render_pdf_page_at_dpi

    match render_pdf_page_at_dpi(source_pdf_path, int(page_data.get("page_num", 0)), export_dpi):
        case Ok(rendered):
            try:
                return _encode_background(rendered)
            finally:
                rendered.close()
        case Err(reason):
            logger.warning("Fondo de alta resolución descartado, se usa el lienzo: %s", reason)
            return None


def build_pptx_export(payload: dict, source_pdf_path: Path | None = None) -> bytes:
    """
    Clona la arquitectura PDF exportándola transparentemente a slides de PowerPoint (.pptx).
    Escala BBoxes usando EMU equivalentes para retener integridad sin importar el DPI de fondo.

    El fondo se re-rasteriza del PDF original a `payload["export_dpi"]` cuando es
    posible: el lienzo trabaja a 100 DPI por velocidad de OCR, y meter ese raster
    en el PPTX es lo que hacía que todo lo no editado se viese blando frente al
    texto editado, que sí es vectorial.
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("La librería python-pptx no se encuentra instalada.")

    assert Presentation is not None
    assert Pt is not None
    assert Emu is not None
    assert RGBColor is not None
    assert PP_ALIGN is not None
    assert MSO_ANCHOR is not None

    prs = Presentation()
    export_dpi = resolve_export_dpi(payload.get("export_dpi"))

    # 914400 EMUs dictaminan 1 pulgada geométrica según la Open XML Specification
    for page_index, page_data in enumerate(payload.get("pages", [])):
        b64_img = page_data.get("image_base64")
        if not b64_img:
            continue
            
        canvas_bytes = _decode_image(b64_img)

        # Invocamos el Slide Master "Vacio"
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Calibrador Dinámico de Presentaciones atadas al Canvas.
        # OJO: las bboxes de los bloques viven en píxeles del lienzo, así que las
        # ratios de mapeo se calculan SIEMPRE con el tamaño del lienzo. Si se
        # tomaran del fondo de alta resolución, cada bloque de texto aterrizaría
        # a una fracción de su posición real.
        from PIL import Image
        with Image.open(io.BytesIO(canvas_bytes)) as canvas_img:
            width_px = float(canvas_img.width)
            height_px = float(canvas_img.height)

        img_stream = io.BytesIO(_high_res_background(page_data, source_pdf_path, export_dpi) or canvas_bytes)
            
        page_w_pt = float(page_data.get("page_width_pt") or width_px)
        page_h_pt = float(page_data.get("page_height_pt") or height_px)
        if page_index == 0:
            # Usar puntos reales del PDF evita depender del DPI de rasterizado.
            prs.slide_width = Emu(int(max(1.0, page_w_pt) * EMU_PER_PT))
            prs.slide_height = Emu(int(max(1.0, page_h_pt) * EMU_PER_PT))

        slide_width_emu = int(cast(Any, prs.slide_width))
        slide_height_emu = int(cast(Any, prs.slide_height))
        
        # Tapizado del Background completo Slide -> Slide
        slide.shapes.add_picture(img_stream, Emu(0), Emu(0), width=Emu(slide_width_emu), height=Emu(slide_height_emu))
        
        # Ecuación de Escalado y Mapeo Posicional del texto
        ratio_x = slide_width_emu / max(1.0, width_px)
        ratio_y = slide_height_emu / max(1.0, height_px)
        _, px_to_pt_y = _page_scale_to_ppt_points(page_data, width_px, height_px)
        
        for block in page_data.get("blocks", []):
            x0, y0, x1, y1 = block["bbox"]
            
            shape_left = Emu(int(x0 * ratio_x))
            shape_top = Emu(int(y0 * ratio_y))
            shape_width = Emu(int((x1 - x0) * ratio_x))
            shape_height = Emu(int((y1 - y0) * ratio_y))
            
            # Recuperar de JSON
            bg_hex = block.get("bg_color", "#ffffff")
            txt_hex = block.get("text_color", "#000000")
            fsize = block.get("font_size", 16)
            font_fam = block.get("font_family", "Arial")

            is_modified = bool(block.get("is_modified", False))
            export_mode = payload.get("export_mode", "only_modified")
            if export_mode == "only_modified" and not is_modified:
                continue

            from pptx.enum.shapes import MSO_SHAPE
            from pptx.util import Inches

            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, shape_left, shape_top, shape_width, shape_height
            )
            if not block.get("bg_transparent", False):
                bg_shape.fill.solid()
                bg_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(bg_hex))
            else:
                bg_shape.fill.background()
                
            # Remover borde azul por default en Powerpoint
            bg_shape.line.fill.background()
            # Remover sombra del rectángulo
            bg_shape.shadow.inherit = False
            text_frame = bg_shape.text_frame
            
            # Cúpula de TextFrame - MEJORADO para alineación perfecta
            text_frame.clear() 
            text_frame.margin_top = Inches(0.02)     # Margen superior mínimo
            text_frame.margin_bottom = Inches(0.02)  # Margen inferior mínimo
            text_frame.margin_left = Inches(0.03)    # Margen izquierdo mínimo
            text_frame.margin_right = Inches(0.03)   # Margen derecho mínimo
            # Mantener el ajuste al ancho del bloque como en el editor canvas.
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = text_frame.paragraphs[0]
            p.text = block.get("text", "")

            # Alineación de párrafo
            align_str = block.get("text_align", "left")
            if align_str == "center":
                p.alignment = PP_ALIGN.CENTER
            elif align_str == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
            
            # Conversión tipográfica exacta según proporción real canvas->página.
            # Se usa el eje Y para preservar altura visual de línea.
            real_pt = max(1.0, float(fsize) * px_to_pt_y)
            p.font.size = Pt(real_pt)
            p.font.name = font_fam
            p.font.color.rgb = RGBColor(*_hex_to_rgb(txt_hex))
            p.font.bold = bool(block.get("is_bold", False))
            p.font.italic = bool(block.get("is_italic", False))
            p.font.underline = bool(block.get("is_underline", False))
            
            # Line spacing para mejor espaciado
            p.line_spacing = _safe_line_spacing(block.get("line_spacing", 1.15))

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def generate_export_zip(payload_dict: dict, source_pdf_path: Path | None = None) -> tuple[Path, Path]:
    """
    Rutea la producción de ambos documentos solicitados al Core
    y los envuelve en un Zip FileResponse para FastAPI.
    """
    temp_dir = Path(tempfile.mkdtemp())
    zip_path = temp_dir / "export.zip"

    raw_targets = payload_dict.get("export_targets") or {}
    export_pdf = bool(raw_targets.get("pdf", True))
    export_pptx = bool(raw_targets.get("pptx", True))
    export_md = bool(raw_targets.get("md", True))

    # Guard rail por robustez: evitar ZIP vacío cuando el cliente no define nada explícito.
    if not (export_pdf or export_pptx or export_md):
        export_pdf = True
        export_pptx = True
        export_md = True
    
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zipf:
        if export_pdf:
            if source_pdf_path and source_pdf_path.exists():
                pdf_bytes = build_pdf_export_from_original(payload_dict, source_pdf_path)
            else:
                pdf_bytes = build_pdf_export(payload_dict)
            zipf.writestr("Presentacion_Editada_Impresa.pdf", pdf_bytes)

        if export_md:
            from .markdown_exporter import build_markdown_export
            zipf.writestr("Presentacion_Editada_DBV.md", build_markdown_export(payload_dict, source_pdf_path))

        if export_pptx and PPTX_AVAILABLE:
            pptx_bytes = build_pptx_export(payload_dict, source_pdf_path)
            zipf.writestr("Presentacion_Editada_DBV.pptx", pptx_bytes)
            
    return zip_path, temp_dir
